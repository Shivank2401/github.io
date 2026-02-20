from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Helper function to add a slide with title and content
    def add_slide(layout_index, title_text):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        return slide

    # Helper to add a placeholder box for images
    def add_image_placeholder(slide, text, left, top, width, height):
        shape = slide.shapes.add_shape(
            1, # msoShapeRectangle
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 220, 220) # Light Gray
        shape.line.color.rgb = RGBColor(100, 100, 100)
        
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Helper for bullet points
    def add_bullet(tf, text, level=0, bold=False, font_size=None):
        p = tf.add_paragraph()
        p.text = text
        p.level = level
        if bold:
            p.font.bold = True
        if font_size:
            p.font.size = Pt(font_size)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Shivank Rajput"
    subtitle.text = "Strategic Analytics & Operations Leader\n\nrajput.shivank34@gmail.com | +91 8318875772\nNoida, IN"

    # 2. Executive Profile
    slide = add_slide(1, "Executive Profile")
    content = slide.placeholders[1].text_frame
    add_bullet(content, "Professional Summary", bold=True)
    add_bullet(content, "Results-focused Deputy Manager with 5+ years of experience. Expert in bridging the gap between raw data and CXO-level strategy.", level=0)
    
    add_bullet(content, "Key Achievements", bold=True)
    add_bullet(content, "â€¢ Inventory: 15% reduction in defects via real-time logic.", level=0)
    add_bullet(content, "â€¢ Sustainability: 20% Carbon Footprint reduction tracked via Python tools.", level=0)
    add_bullet(content, "â€¢ Growth: 25% Sales Conversion boost via city-level data mining.", level=0)

    # 3. Competency Matrix (Enhanced Version)
    # Changed to Layout 5 (Title Only) to avoid XML manipulation errors
    slide = add_slide(5, "Competency Matrix")

    # Layout dimensions
    top_headers = Inches(1.5)
    top_content = Inches(2.0)
    width = Inches(3.2)
    height = Inches(5.0)
    font_size = 12

    # Column 1: Technical & Analytical
    shapes = slide.shapes
    
    # Header 1
    tb1 = shapes.add_textbox(Inches(0.3), top_headers, width, Inches(0.5))
    p = tb1.text_frame.add_paragraph()
    p.text = "ðŸ”¹ Technical & Analytical"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 112, 192) # Blue

    # Content 1
    tb1_c = shapes.add_textbox(Inches(0.3), top_content, width, height)
    tf = tb1_c.text_frame
    tf.word_wrap = True
    skills_tech = [
        "Data Analysis & Mining",
        "MIS Reporting Architecture",
        "Reporting Automation",
        "Process Optimization",
        "Primary & Secondary Research",
        "Insight Generation"
    ]
    for skill in skills_tech:
        add_bullet(tf, f"â€¢ {skill}", font_size=font_size)

    # Column 2: Domain Expertise
    # Header 2
    tb2 = shapes.add_textbox(Inches(3.6), top_headers, width, Inches(0.5))
    p = tb2.text_frame.add_paragraph()
    p.text = "ðŸ”¹ Domain Expertise"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 112, 192)

    # Content 2
    tb2_c = shapes.add_textbox(Inches(3.6), top_content, width, height)
    tf = tb2_c.text_frame
    tf.word_wrap = True
    skills_domain = [
        "Amazon/Flipkart/Walmart APIs",
        "Retail & E-commerce Ops",
        "CXO Board Deck Creation",
        "Project Management",
        "Ops Workflow Management",
        "Team Management & Training",
        "Cross-functional Collaboration"
    ]
    for skill in skills_domain:
        add_bullet(tf, f"â€¢ {skill}", font_size=font_size)

    # Column 3: Tools
    # Header 3
    tb3 = shapes.add_textbox(Inches(6.9), top_headers, width, Inches(0.5))
    p = tb3.text_frame.add_paragraph()
    p.text = "ðŸ”¹ Tools & Stack"
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 112, 192)

    # Content 3
    tb3_c = shapes.add_textbox(Inches(6.9), top_content, width, height)
    tf = tb3_c.text_frame
    tf.word_wrap = True
    skills_tools = [
        "Python (ETL & Automation)",
        "Power BI & Tableau",
        "Azure & SQL",
        "Power Automate",
        "VS Code & Git",
        "Figma & Canva"
    ]
    for skill in skills_tools:
        add_bullet(tf, f"â€¢ {skill}", font_size=font_size)

    # 4. Enterprise Ecosystem
    slide = add_slide(1, "Enterprise Ecosystem: Internal Portals")
    content = slide.placeholders[1].text_frame
    content.text = "Architected and managed critical business tools:"
    
    add_bullet(content, "â€¢ Datahive: Centralized data warehouse (Azure + Power BI).")
    add_bullet(content, "â€¢ HRMS Suite: Admin & Employee portals for attrition analysis.")
    add_bullet(content, "â€¢ Ops Trackers: Project, Query, and Asset Trackers (Digitized Workflows).")
    
    add_image_placeholder(slide, "INSERT SCREENSHOT: Datahive or HRMS Interface", Inches(1), Inches(4.5), Inches(8), Inches(2.5))

    # 5. Automation Workflow
    slide = add_slide(1, "Automation & Operational Excellence")
    content = slide.placeholders[1].text_frame
    content.text = "The 'Operational Tool' Logic (Python + Power Automate)"
    
    # Draw a visual flow using text boxes
    shapes = slide.shapes
    y_pos = Inches(3)
    
    # Step 1
    shapes.add_textbox(Inches(0.5), y_pos, Inches(2), Inches(1)).text_frame.text = "1. INPUT\nRaw Amazon/Walmart Files"
    # Arrow
    shapes.add_shape(33, Inches(2.6), y_pos+Inches(0.2), Inches(0.5), Inches(0.5)) # Right Arrow
    
    # Step 2
    shapes.add_textbox(Inches(3.2), y_pos, Inches(2), Inches(1)).text_frame.text = "2. PROCESS\nPython ETL (VS Code)"
    # Arrow
    shapes.add_shape(33, Inches(5.3), y_pos+Inches(0.2), Inches(0.5), Inches(0.5))
    
    # Step 3
    shapes.add_textbox(Inches(5.9), y_pos, Inches(2), Inches(1)).text_frame.text = "3. TRIGGER\nPower Automate Approvals"
    # Arrow
    shapes.add_shape(33, Inches(8.0), y_pos+Inches(0.2), Inches(0.5), Inches(0.5))
    
    # Step 4
    shapes.add_textbox(Inches(8.6), y_pos, Inches(2), Inches(1)).text_frame.text = "4. VISUALIZE\nLive Power BI Dashboard"

    add_image_placeholder(slide, "INSERT SCREENSHOT: Python Script or Power Automate Flow", Inches(2), Inches(5), Inches(6), Inches(2))

    # 6. Power BI Portfolio
    slide = add_slide(1, "Advanced Visualization Portfolio")
    content = slide.placeholders[1].text_frame
    content.text = "Strategic Dashboards"
    
    add_bullet(content, "Logistics & Ops: Inventory, Aging, Amazon Shipping, Demand vs Sales.")
    add_bullet(content, "Sales & Strategy: P&L Analysis, Customer Buying Patterns, Returns Analysis.")
    add_bullet(content, "Design Philosophy: Prototyping in Figma/Canva before build.")

    add_image_placeholder(slide, "INSERT SCREENSHOT: Collage of Dashboards", Inches(1), Inches(4), Inches(8), Inches(3))

    # 7. Contact
    slide = add_slide(1, "Contact")
    content = slide.placeholders[1].text_frame
    content.text = "Ready to Transform Your Data?"
    
    add_bullet(content, "Shivank Rajput", bold=True)
    add_bullet(content, "rajput.shivank34@gmail.com")
    add_bullet(content, "+91 8318875772")
    add_bullet(content, "Noida, IN")

    prs.save('Shivank_Rajput_Portfolio_Updated.pptx')

create_presentation()